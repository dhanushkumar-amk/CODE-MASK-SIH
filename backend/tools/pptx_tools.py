"""Real PowerPoint-generation tool, sandboxed to backend/workspace/.

Produces actual .pptx files via python-pptx so the agent can deliver
presentations alongside Word documents and text files.

Expected presentation structure (see pptx_generate):
    Slide 0           -> title slide (title as main title)
    Slides 1..N       -> "Title and Content" layout, one per entry in
                         `slides` (each a {"heading", "body"} dict)

Future extension points (not implemented yet):
    - bullets: accept a "bullets" list per slide for multi-point bodies
    - speaker notes, images, charts via extra per-slide keys
The slide-building loop is isolated so those additions only extend it.
"""

import re
import sys
from pathlib import Path

# Make `tools` importable both when this module is imported by the app
# (backend on sys.path) and when run directly (backend/tools on sys.path).
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from tools.file_tools import WORKSPACE_DIR, _resolve_safe  # noqa: E402

try:
    from pptx import Presentation
except ImportError:
    # python-pptx missing: pptx_generate reports a clean error at call
    # time instead of failing the whole module import.
    Presentation = None


def _slugify(title: str) -> str:
    """Turn a title into a safe default filename: lowercase, underscores."""
    slug = title.lower().strip()
    slug = re.sub(r"[^a-z0-9]+", "_", slug)
    slug = slug.strip("_")
    return slug or "presentation"


def _validate_slides(slides) -> str | None:
    """Return an error message if `slides` is malformed, else None.

    Each slide must be a dict with non-empty "heading" and "body" string
    keys. The model sometimes emits strings or missing keys, and a broken
    slide list must produce a clean error, not a broken file.
    """
    if not isinstance(slides, list) or not slides:
        return "slides must be a non-empty list of slide dicts."

    for index, slide in enumerate(slides, 1):
        if not isinstance(slide, dict):
            return f"slide {index} is not a dict: {slide!r}"
        heading = slide.get("heading")
        body = slide.get("body")
        if not isinstance(heading, str) or not heading.strip():
            return f"slide {index} is missing a non-empty 'heading'."
        if not isinstance(body, str) or not body.strip():
            return f"slide {index} is missing a non-empty 'body'."
    return None


def pptx_generate(
    title: str,
    slides: list,
    filename: str = None,
    provenance_record: dict = None,
) -> dict:
    """Create a PowerPoint presentation in the workspace and save it.

    Args:
        title: Presentation title, rendered on a title slide.
        slides: List of dicts, each {"heading": "...", "body": "..."}.
        filename: Optional output name. Defaults to a slug of the title
            (e.g. "Weekly Review" -> "weekly_review.pptx"). ".pptx" is
            appended if missing.
        provenance_record: Optional provenance metadata dict. When provided,
            adds a final "Provenance & Audit Trail" slide.

    Returns:
        {"status": "success", "output": "Presentation saved as <name>",
         "file_path": "<full path>"} on success, or
        {"status": "error", "output": "<clear error message>"} on failure.
        Never raises.
    """
    if Presentation is None:
        return {
            "status": "error",
            "output": "python-pptx is not installed. Run: pip install python-pptx",
        }

    error = _validate_slides(slides)
    if error:
        return {"status": "error", "output": error}

    if not isinstance(title, str) or not title.strip():
        return {"status": "error", "output": "Title must be a non-empty string."}

    if filename is None or not str(filename).strip():
        filename = _slugify(title)
    else:
        filename = str(filename)

    if not filename.lower().endswith(".pptx"):
        filename += ".pptx"

    # Reuse file_tools' sandbox: same workspace root, same traversal
    # rejection, so the deliverable lands next to the agent's other files.
    safe_path = _resolve_safe(filename)
    if safe_path is None:
        return {
            "status": "error",
            "output": f"Path rejected (must stay inside the workspace): {filename!r}",
        }

    try:
        prs = Presentation()

        # Title slide: layout 0 of the default template.
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title

        # One "Title and Content" slide per entry.
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data["heading"]
            slide.placeholders[1].text = slide_data["body"]

        # Append Provenance & Audit Trail slide if requested or available
        try:
            from provenance.provenance_tracker import (
                format_provenance_footer,
                get_current_provenance_record,
            )
            if provenance_record is None and get_current_provenance_record is not None:
                provenance_record = get_current_provenance_record()

            if provenance_record and format_provenance_footer is not None:
                prov_slide = prs.slides.add_slide(prs.slide_layouts[1])
                prov_slide.shapes.title.text = "Provenance & Audit Trail"
                footer_text = format_provenance_footer(provenance_record)
                prov_slide.placeholders[1].text = footer_text
        except Exception as prov_exc:
            print(f"[PPTX] Provenance slide skipped: {prov_exc}")

        safe_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(safe_path))
    except OSError as exc:
        return {"status": "error", "output": f"Could not save presentation: {exc}"}
    except Exception as exc:  # noqa: BLE001 - library errors must not crash.
        return {"status": "error", "output": f"pptx_generate failed: {exc}"}

    return {
        "status": "success",
        "output": f"Presentation saved as {filename}",
        "file_path": str(safe_path),
    }


if __name__ == "__main__":
    TITLE = "Weekly Inspection Summary"
    SLIDES = [
        {
            "heading": "Findings",
            "body": "Minor corrosion observed on joint A-12.",
        },
        {
            "heading": "Recommendation",
            "body": "Schedule maintenance within 30 days.",
        },
        {
            "heading": "Status",
            "body": "No immediate safety risk identified.",
        },
    ]

    print(f"Workspace: {WORKSPACE_DIR}\n")

    result = pptx_generate(TITLE, SLIDES)
    print("RESULT:", result)

    if result["status"] == "success":
        path = Path(result["file_path"])
        print(f"\nFile exists: {path.exists()} | Size: {path.stat().st_size} bytes")

        # Read the file back with python-pptx to confirm it is a real,
        # openable presentation with the right slide count and content.
        from pptx import Presentation as Prs

        prs = Prs(str(path))
        print(f"Slide count: {len(prs.slides)} (expected 1 title + {len(SLIDES)})")
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
            print(f"  slide {i}: {texts}")
