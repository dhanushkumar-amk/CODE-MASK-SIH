import sys
import os
import zipfile
from pathlib import Path
from docx import Document
from pptx import Presentation as Prs
from openpyxl import load_workbook

# Put backend root on sys.path
SYS_BACKEND_DIR = Path(__file__).resolve().parent
if str(SYS_BACKEND_DIR) not in sys.path:
    sys.path.insert(0, str(SYS_BACKEND_DIR))

from tools.file_tools import WORKSPACE_DIR
from tools.docx_tools import docx_generate
from tools.pptx_tools import pptx_generate
from tools.xlsx_tools import xlsx_generate
from rag.vector_store import add_documents, query_knowledge_base
from ocr.ocr_tools import extract_text_from_image
from provenance.provenance_tracker import (
    reset_provenance_context,
    set_routing_context,
    add_source_document,
    get_current_provenance_record,
    format_provenance_footer,
)

print("=" * 60)
print(" END-TO-END PROVENANCE & AUDIT TRAIL TEST SUITE")
print("=" * 60)

# Create mock scan image and SOP text file in workspace for real SHA-256 hashing
scan_file = WORKSPACE_DIR / "pipe_scan.png"
sop_file = WORKSPACE_DIR / "SOP_Valve_Maintenance_Rev3.txt"

with open(sop_file, "w", encoding="utf-8") as f:
    f.write("SOP Rev3: Valves must be inspected every 90 days. Joint corrosion requires clean & coat within 30 days.\n")

# Create a small dummy png
from PIL import Image
img = Image.new("RGB", (100, 50), color="white")
img.save(scan_file)

# Ingest SOP chunk into RAG
add_documents(
    chunks=["Valves must be inspected every 90 days."],
    metadatas=[{"source": "SOP_Valve_Maintenance_Rev3.txt"}],
    ids=["sop_valve_rev3"]
)

# ---------------------------------------------------------------------------
# Test 1: Full "scan -> OCR -> RAG lookup -> draft approval note -> docx"
# ---------------------------------------------------------------------------
print("\n[TEST 1] Full RAG + Scan End-to-End Docx Generation")
reset_provenance_context()
set_routing_context("qwen2.5:1.5b-instruct", "document")

# Step 1: Simulate OCR scan read
print("  1. Simulating OCR on pipe_scan.png...")
add_source_document("pipe_scan.png")

# Step 2: Simulate RAG retrieval
print("  2. Simulating RAG retrieval for 'valve maintenance'...")
rag_res = query_knowledge_base("valve maintenance", n_results=1)
print(f"     Retrieved sources: {rag_res.get('sources')}")

# Step 3: Generate Docx deliverable
print("  3. Generating DOCX deliverable...")
docx_res = docx_generate(
    title="Pipe Inspection Approval Note",
    content="Findings: Minor corrosion observed on joint A-12.\n\nRecommendation: Clean, coat, and schedule re-inspection within 30 days as per SOP."
)
print(f"     Docx result: {docx_res['status']} | {docx_res['output']}")

# Inspect generated DOCX
docx_path = Path(docx_res["file_path"])
doc = Document(str(docx_path))
docx_text = "\n".join(p.text for p in doc.paragraphs)

print("\n--- DOCX PROVENANCE SECTION TEXT DUMP ---")
print(docx_text[docx_text.find("Provenance & Audit Trail"):])
print("-----------------------------------------")

assert "Provenance & Audit Trail" in docx_text, "Heading 3 'Provenance & Audit Trail' missing from DOCX"
assert "qwen2.5:1.5b-instruct (local, offline)" in docx_text, "Model string missing from DOCX"
assert "pipe_scan.png (SHA-256:" in docx_text, "pipe_scan.png SHA-256 hash missing from DOCX"
assert "SOP_Valve_Maintenance_Rev3.txt (SHA-256:" in docx_text, "SOP SHA-256 hash missing from DOCX"
assert "Verified offline" in docx_text, "Network status missing from DOCX"
print("[OK] TEST 1 PASSED: Real SHA-256 source document fingerprints verified in .docx!")

# ---------------------------------------------------------------------------
# Test 2: No-RAG Source Scenario (General Model Knowledge)
# ---------------------------------------------------------------------------
print("\n[TEST 2] No-RAG Source Scenario (General Model Knowledge)")
reset_provenance_context()
set_routing_context("qwen2.5:1.5b-instruct", "coding")

norag_docx = docx_generate(
    title="Duck Number Algorithm",
    content="A Duck number is a number which has zero present in it, but not at the beginning."
)

doc_norag = Document(str(norag_docx["file_path"]))
text_norag = "\n".join(p.text for p in doc_norag.paragraphs)

print("\n--- NO-RAG PROVENANCE SECTION TEXT DUMP ---")
print(text_norag[text_norag.find("Provenance & Audit Trail"):])
print("-------------------------------------------")

assert "Provenance & Audit Trail" in text_norag, "Heading missing from No-RAG DOCX"
assert "No external source documents used — general model knowledge" in text_norag, "No-RAG fallback string missing"
print("[OK] TEST 2 PASSED: No-RAG fallback handled gracefully!")

# ---------------------------------------------------------------------------
# Test 3: PPTX Provenance Slide Test
# ---------------------------------------------------------------------------
print("\n[TEST 3] PPTX Presentation Final Slide Test")
reset_provenance_context()
set_routing_context("qwen2.5:1.5b-instruct", "document")
add_source_document("SOP_Valve_Maintenance_Rev3.txt")

pptx_res = pptx_generate(
    title="Pipeline Review",
    slides=[{"heading": "Status", "body": "All valves operational."}]
)

prs = Prs(str(pptx_res["file_path"]))
last_slide = prs.slides[-1]
last_slide_texts = [s.text_frame.text for s in last_slide.shapes if s.has_text_frame]

print("\n--- PPTX PROVENANCE SLIDE TEXT DUMP ---")
for t in last_slide_texts:
    print(t)
print("---------------------------------------")

assert any("Provenance & Audit Trail" in t for t in last_slide_texts), "Provenance slide title missing"
assert any("SOP_Valve_Maintenance_Rev3.txt (SHA-256:" in t for t in last_slide_texts), "PPTX SHA-256 missing"
print("[OK] TEST 3 PASSED: PPTX Provenance slide created!")

# ---------------------------------------------------------------------------
# Test 4: XLSX Provenance Worksheet Tab Test
# ---------------------------------------------------------------------------
print("\n[TEST 4] XLSX Spreadsheet Provenance Tab Test")
reset_provenance_context()
set_routing_context("qwen2.5:1.5b-instruct", "document")
add_source_document("SOP_Valve_Maintenance_Rev3.txt")

xlsx_res = xlsx_generate(
    title="Inspection Log",
    headers=["Valve ID", "Status"],
    rows=[["V-101", "Passed"]]
)

wb = load_workbook(str(xlsx_res["file_path"]))
assert "Provenance" in wb.sheetnames, "Provenance worksheet tab missing in XLSX"

prov_tab = wb["Provenance"]
prov_rows = [list(r) for r in prov_tab.iter_rows(values_only=True)]

print("\n--- XLSX PROVENANCE TAB ROWS ---")
for r in prov_rows:
    print(f"  {r[0]:<20}: {r[1]}")
print("--------------------------------")

assert prov_rows[0] == ["Field", "Audit Trail Detail"]
assert prov_rows[1][1] == "qwen2.5:1.5b-instruct (local, offline)"
assert "SOP_Valve_Maintenance_Rev3.txt (SHA-256:" in prov_rows[4][1]
print("[OK] TEST 4 PASSED: XLSX Provenance worksheet tab created!")

print("\n" + "=" * 60)
print(" ALL 4 PROVENANCE AUDIT TRAIL TESTS PASSED PERFECTLY!")
print("=" * 60)
